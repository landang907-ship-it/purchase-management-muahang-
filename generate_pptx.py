import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    # Set slide width and height to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Color Palette
    DARK_BG = RGBColor(15, 23, 42)      # Slate 900
    TEXT_DARK = RGBColor(30, 41, 59)    # Slate 800
    TEXT_LIGHT = RGBColor(248, 250, 252)# Slate 50
    TEXT_MUTED = RGBColor(100, 116, 139)# Slate 500
    PRIMARY = RGBColor(37, 99, 235)     # Blue 600
    PRIMARY_LIGHT = RGBColor(239, 246, 255) # Blue 50
    CARD_BG = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(226, 232, 240)
    ACCENT_GREEN = RGBColor(16, 185, 129) # Emerald 500
    ACCENT_ORANGE = RGBColor(245, 158, 11) # Amber 500

    def add_header(slide, title_vi, title_zh, category_text="DỰ ÁN QUẢN LÝ MUA HÀNG • 采购管理系统"):
        # Header banner category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY

        # Header Title (Bilingual)
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.95))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        
        p_vi = tf_title.paragraphs[0]
        p_vi.text = title_vi
        p_vi.font.size = Pt(21)
        p_vi.font.bold = True
        p_vi.font.color.rgb = TEXT_DARK
        p_vi.space_after = Pt(2)

        p_zh = tf_title.add_paragraph()
        p_zh.text = title_zh
        p_zh.font.size = Pt(14)
        p_zh.font.bold = True
        p_zh.font.color.rgb = PRIMARY

    def create_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # ==================== SLIDE 1: Title Slide (Dark Theme - Bilingual) ====================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_BG
    bg1.line.fill.background()

    card_dec = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.1))
    card_dec.fill.solid()
    card_dec.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate 800
    card_dec.line.color.rgb = PRIMARY
    card_dec.line.width = Pt(2)

    tb1 = slide1.shapes.add_textbox(Inches(1.4), Inches(1.5), Inches(10.5), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "BÁO CÁO TÍNH NĂNG DỰ ÁN • 项目功能报告 (越中双语版)"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p1.space_after = Pt(10)

    p2 = tf1.add_paragraph()
    p2.text = "HỆ THỐNG QUẢN LÝ MUA HÀNG"
    p2.font.size = Pt(30)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_LIGHT

    p2_zh = tf1.add_paragraph()
    p2_zh.text = "采购管理系统 (Purchase Management System)"
    p2_zh.font.size = Pt(20)
    p2_zh.font.bold = True
    p2_zh.font.color.rgb = RGBColor(147, 197, 253) # Light Blue
    p2_zh.space_after = Pt(15)

    p3 = tf1.add_paragraph()
    p3.text = "Giải pháp số hóa quy trình mua sắm vật tư doanh nghiệp & nhà máy sản xuất"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(226, 232, 240)

    p3_zh = tf1.add_paragraph()
    p3_zh.text = "企业与生产工厂物料采购全流程数字化解决方案"
    p3_zh.font.size = Pt(13)
    p3_zh.font.color.rgb = RGBColor(148, 163, 184)
    p3_zh.space_after = Pt(25)

    p4 = tf1.add_paragraph()
    p4.text = "📅 Cập nhật / 更新: 07/2026  |  🌐 Ngôn ngữ / 语言: Việt - Trung (越中)  |  💻 React 19 + Vite 6 + Supabase"
    p4.font.size = Pt(11)
    p4.font.color.rgb = RGBColor(203, 213, 225)

    # ==================== SLIDE 2: Tổng Quan & Mục Tiêu ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "TỔNG QUAN HỆ THỐNG & MỤC TIÊU DỰ ÁN", "系统概述与项目目标")

    card_width = Inches(3.64)
    card_height = Inches(4.9)
    top_pos = Inches(1.8)

    features_s2 = [
        {
            "title_vi": "🎯 Mục Tiêu Chính",
            "title_zh": "主要目标",
            "subtitle": "Số hóa quy trình mua sắm / 采购数字化",
            "items": [
                ("Quản lý tập trung toàn bộ đề xuất đơn mua", "集中管理所有采购申请与订单"),
                ("Tối ưu thời gian duyệt giữa các phòng ban", "优化部门间审批处理时间"),
                ("Giảm thiểu sai sót và trùng lặp mã vật tư", "减少物料编码错误与重复"),
                ("Theo dõi thời gian thực tiến độ xử lý đơn", "实时跟踪订单处理进度")
            ]
        },
        {
            "title_vi": "🏭 Môi Trường Áp Dụng",
            "title_zh": "适用环境",
            "subtitle": "Đặc thù nhà máy / 工厂与企业特性",
            "items": [
                ("Hỗ trợ Đa ngôn ngữ (Tiếng Việt & Tiếng Trung)", "支持多语言（越南语与中文）"),
                ("Phù hợp doanh nghiệp sản xuất / liên doanh", "适用于生产型企业/合资企业"),
                ("Quản lý theo mô hình Phân xưởng / Bộ phận", "按车间/部门架构层级管理"),
                ("Phân quyền minh bạch Nhân viên & Admin", "员工与管理员权限透明划分")
            ]
        },
        {
            "title_vi": "📱 Trải Nghiệm Đa Nền Tảng",
            "title_zh": "多平台体验",
            "subtitle": "Desktop & Mobile Responsive",
            "items": [
                ("Desktop DataTable hiển thị đa thông tin", "桌面端多信息数据表界面"),
                ("Mobile Card List tiện thao tác cảm ứng", "移动端卡片列表触摸优化"),
                ("Tối ưu tốc độ tải trang với Vite 6", "Vite 6 极致页面加载速度"),
                ("Thông báo Toast tức thì cho mọi thao tác", "操作实时 Toast 消息通知")
            ]
        }
    ]

    for i, f in enumerate(features_s2):
        left_pos = Inches(0.8 + i * 4.04)
        create_card(slide2, left_pos, top_pos, card_width, card_height)
        
        tb = slide2.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(0.15), card_width - Inches(0.3), card_height - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f["title_vi"] + " • " + f["title_zh"]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        p_sub = tf.add_paragraph()
        p_sub.text = f["subtitle"]
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = TEXT_MUTED
        p_sub.space_after = Pt(10)

        for vi, zh in f["items"]:
            pi_vi = tf.add_paragraph()
            pi_vi.text = "• " + vi
            pi_vi.font.size = Pt(11)
            pi_vi.font.bold = True
            pi_vi.font.color.rgb = TEXT_DARK
            
            pi_zh = tf.add_paragraph()
            pi_zh.text = "  " + zh
            pi_zh.font.size = Pt(10)
            pi_zh.font.color.rgb = TEXT_MUTED
            pi_zh.space_after = Pt(6)

    # ==================== SLIDE 3: Kiến Trúc & Công Nghệ ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "KIẾN TRÚC CÔNG NGHỆ (TECH STACK)", "技术架构与代码结构")

    create_card(slide3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9))
    tb_l = slide3.shapes.add_textbox(Inches(0.95), Inches(1.95), Inches(5.3), Inches(4.6))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "💻 Công Nghệ (Tech Stack / 现代技术栈)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(10)

    stack_items = [
        ("React 19 & Vite 6", "UI Framework mới nhất, build tức thì / 最新UI框架，即时构建"),
        ("TypeScript 5.8", "Type-safety, giảm lỗi runtime / 类型安全，减少运行时错误"),
        ("Tailwind CSS v4", "Styling hiện đại, chuẩn UI/UX / 现代样式，UI/UX 标准"),
        ("Supabase (Postgres)", "Auth, Database Realtime & Storage / 实时数据库与存储"),
        ("Lucide & Motion", "Icon chuyên nghiệp & hiệu ứng mượt / 专业图标与流畅动画")
    ]

    for title, desc in stack_items:
        p1 = tf_l.add_paragraph()
        p1.text = "✔ " + title
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = TEXT_DARK
        
        p2 = tf_l.add_paragraph()
        p2.text = "   " + desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_after = Pt(6)

    create_card(slide3, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.9), bg_color=PRIMARY_LIGHT, border_color=PRIMARY)
    tb_r = slide3.shapes.add_textbox(Inches(6.95), Inches(1.95), Inches(5.4), Inches(4.6))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "🏗️ Cấu Trúc (Feature-Sliced Design / 代码结构)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(10)

    fsd_items = [
        ("features/auth", "Đăng nhập, phân quyền Role, Profile & User Admin\n登录、角色权限、个人资料与用户管理"),
        ("features/purchase", "Đơn mua hàng, Bộ lọc, Mã vật tư, Modal, Excel\n采购订单管理、筛选、物料编码、详情弹窗、Excel"),
        ("features/notifications", "Thông báo biến động đơn hàng real-time\n订单变动实时通知系统"),
        ("shared & layout", "UI Component dùng chung (Toast, ProtectedRoute)\n通用 UI 组件")
    ]

    for folder, desc in fsd_items:
        p1 = tf_r.add_paragraph()
        p1.text = "📁 " + folder
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = PRIMARY
        
        p2 = tf_r.add_paragraph()
        p2.text = "   " + desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_DARK
        p2.space_after = Pt(8)

    # ==================== SLIDE 4: Tính Năng Đơn Mua Hàng ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "TÍNH NĂNG 1: QUẢN LÝ ĐƠN MUA HÀNG HỆ THỐNG", "功能 1: 系统采购订单管理")

    grid_w = Inches(5.6)
    grid_h = Inches(2.3)

    coords = [
        (Inches(0.8), Inches(1.8)),
        (Inches(6.8), Inches(1.8)),
        (Inches(0.8), Inches(4.35)),
        (Inches(6.8), Inches(4.35))
    ]

    cards_data_s4 = [
        {
            "title_vi": "📋 Chế Độ Hiển Thị Linh Hoạt",
            "title_zh": "灵活显示模式",
            "desc_vi": "DataTable đa thông tin trên Desktop và Mobile Purchase List dạng card cảm ứng trên smartphone.",
            "desc_zh": "桌面端多信息数据表与移动端智能卡片列表。"
        },
        {
            "title_vi": "🔍 Bộ Lọc Đa Tiêu Chí Thông Minh",
            "title_zh": "智能多标准筛选",
            "desc_vi": "Lọc theo Xưởng (Workshop), Người yêu cầu (Requester), Trạng thái đơn, Khoảng thời gian & Tìm nhanh.",
            "desc_zh": "按车间、申请人、订单状态、日期范围筛选及快速搜索。"
        },
        {
            "title_vi": "🔍 Xem & Cập Nhật Chi Tiết Đơn",
            "title_zh": "查看与更新订单详情",
            "desc_vi": "PurchaseDetailModal hiển thị chi tiết số lượng, đơn vị, nhà cung cấp, quy cách & lịch sử cập nhật.",
            "desc_zh": "详情弹窗显示数量、单位、供应商、规格及变更历史。"
        },
        {
            "title_vi": "🔄 Chuyển Đổi Trạng Thái Đơn Hàng",
            "title_zh": "订单状态流转",
            "desc_vi": "Chờ duyệt (待审核 03) ➔ Đã duyệt (已批准 05) ➔ Từ chối (已拒绝 08) ➔ Đã tạo PO ➔ Hoàn thành.",
            "desc_zh": "待审核 (03) ➔ 已批准 (05) ➔ 已拒绝 (08) ➔ 已生成PO ➔ 已完成。"
        }
    ]

    for i, cdata in enumerate(cards_data_s4):
        left, top = coords[i]
        create_card(slide4, left, top, grid_w, grid_h)
        
        tb = slide4.shapes.add_textbox(left + Inches(0.15), top + Inches(0.12), grid_w - Inches(0.3), grid_h - Inches(0.24))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = cdata["title_vi"] + " • " + cdata["title_zh"]
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.space_after = Pt(4)

        p_vi = tf.add_paragraph()
        p_vi.text = cdata["desc_vi"]
        p_vi.font.size = Pt(11)
        p_vi.font.color.rgb = TEXT_DARK

        p_zh = tf.add_paragraph()
        p_zh.text = cdata["desc_zh"]
        p_zh.font.size = Pt(10)
        p_zh.font.color.rgb = TEXT_MUTED

    # ==================== SLIDE 5: Dashboard & Analytics ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "TÍNH NĂNG 2: DASHBOARD THỐNG KÊ & BÁO CÁO", "功能 2: 统计仪表板与报告")

    kpis = [
        ("TỔNG SỐ ĐƠN / 订单总数", "1,250", PRIMARY),
        ("ĐƠN CHỜ DUYỆT / 待审核", "45", ACCENT_ORANGE),
        ("ĐƠN ĐÃ DUYỆT / 已批准", "1,120", ACCENT_GREEN),
        ("HOÀN THÀNH / 完成率", "94.8%", PRIMARY)
    ]

    kpi_w = Inches(2.7)
    kpi_h = Inches(1.3)
    for i, (kpi_t, kpi_v, color) in enumerate(kpis):
        left = Inches(0.8 + i * 3.0)
        create_card(slide5, left, Inches(1.8), kpi_w, kpi_h, bg_color=PRIMARY_LIGHT, border_color=color)
        
        tb = slide5.shapes.add_textbox(left + Inches(0.1), Inches(1.88), kpi_w - Inches(0.2), kpi_h - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = kpi_t
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED

        p_val = tf.add_paragraph()
        p_val.text = kpi_v
        p_val.font.size = Pt(22)
        p_val.font.bold = True
        p_val.font.color.rgb = color

    create_card(slide5, Inches(0.8), Inches(3.3), Inches(5.7), Inches(3.4))
    tb_b1 = slide5.shapes.add_textbox(Inches(0.95), Inches(3.45), Inches(5.4), Inches(3.1))
    tf_b1 = tb_b1.text_frame
    tf_b1.word_wrap = True
    
    p = tf_b1.paragraphs[0]
    p.text = "📈 Biểu Đồ Thống Kê (Recharts / 数据图表)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(8)

    b1_items = [
        ("Trực quan hóa lượng mua hàng theo từng xưởng", "按车间可视化采购申请量"),
        ("Phân tích xu hướng đăng ký mua theo tháng", "月度采购申请趋势分析"),
        ("Theo dõi tỷ lệ phê duyệt & lý do từ chối đơn", "跟踪批准率与拒绝原因")
    ]
    for vi, zh in b1_items:
        pi_vi = tf_b1.add_paragraph()
        pi_vi.text = "✔ " + vi
        pi_vi.font.size = Pt(11)
        pi_vi.font.bold = True
        pi_vi.font.color.rgb = TEXT_DARK
        
        pi_zh = tf_b1.add_paragraph()
        pi_zh.text = "   " + zh
        pi_zh.font.size = Pt(10)
        pi_zh.font.color.rgb = TEXT_MUTED
        pi_zh.space_after = Pt(6)

    create_card(slide5, Inches(6.8), Inches(3.3), Inches(5.7), Inches(3.4))
    tb_b2 = slide5.shapes.add_textbox(Inches(6.95), Inches(3.45), Inches(5.4), Inches(3.1))
    tf_b2 = tb_b2.text_frame
    tf_b2.word_wrap = True
    
    p = tf_b2.paragraphs[0]
    p.text = "💡 Welcome Guide (欢迎指南与体验)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(8)

    b2_items = [
        ("WelcomeGuide hướng dẫn thao tác người dùng mới", "新用户操作指引指南"),
        ("Bảng điều khiển tổng quan các sự kiện sắp tới", "即将到来事件概览面板"),
        ("Chỉ số trực tiếp kết nối dữ liệu Supabase", "直接连接 Supabase 实时数据")
    ]
    for vi, zh in b2_items:
        pi_vi = tf_b2.add_paragraph()
        pi_vi.text = "✔ " + vi
        pi_vi.font.size = Pt(11)
        pi_vi.font.bold = True
        pi_vi.font.color.rgb = TEXT_DARK
        
        pi_zh = tf_b2.add_paragraph()
        pi_zh.text = "   " + zh
        pi_zh.font.size = Pt(10)
        pi_zh.font.color.rgb = TEXT_MUTED
        pi_zh.space_after = Pt(6)

    # ==================== SLIDE 6: Đơn Xử Lý & Mã Vật Tư ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "TÍNH NĂNG 3: ĐƠN ĐÃ XỬ LÝ & QUẢN LÝ MÃ VẬT TƯ", "功能 3: 已处理订单与物料编码管理")

    create_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9))
    tb_p = slide6.shapes.add_textbox(Inches(0.95), Inches(1.95), Inches(5.3), Inches(4.6))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True

    p = tf_p.paragraphs[0]
    p.text = "📦 Đơn Đã Xử Lý (Processed Orders / 已处理订单)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(10)

    p_items = [
        ("Lưu trữ lịch sử đơn hàng:", "Lưu trữ các đơn đã nghiệm thu, giao hàng hoặc hoàn thành.\n归档已完成及已交付的订单历史。"),
        ("Phục vụ đối soát kế toán:", "Dễ dàng tra cứu chứng từ, hóa đơn & đối soát tài chính.\n便于财务与会计进行对账与单据查询。"),
        ("Lọc & Tìm kiếm đơn cũ:", "Tra cứu nhanh theo mã đơn hoặc khoảng ngày.\n按订单号或日期快速检索历史订单。")
    ]
    for title, desc in p_items:
        p1 = tf_p.add_paragraph()
        p1.text = "🔹 " + title
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = TEXT_DARK
        p2 = tf_p.add_paragraph()
        p2.text = "   " + desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_after = Pt(8)

    create_card(slide6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.9))
    tb_m = slide6.shapes.add_textbox(Inches(6.95), Inches(1.95), Inches(5.4), Inches(4.6))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True

    p = tf_m.paragraphs[0]
    p.text = "🏷️ Mã Vật Tư (Material Codes / 物料编码)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(10)

    m_items = [
        ("Danh mục mã SAP & Nội bộ:", "Đồng bộ và quản lý chuẩn hóa mã vật tư, linh kiện.\n同步与标准化管理 SAP 编码与内部编码。"),
        ("Quy chuẩn thông số kỹ thuật:", "Lưu giữ tên gọi, thông số & đơn vị tính chuẩn.\n规范物料名称、规格参数及标准单位。"),
        ("Tránh trùng lặp mã vật tư:", "Giúp nhân viên chọn đúng mã vật tư khi tạo đơn.\n帮助员工在创建订单时准确选择编码。")
    ]
    for title, desc in m_items:
        p1 = tf_m.add_paragraph()
        p1.text = "🔹 " + title
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = TEXT_DARK
        p2 = tf_m.add_paragraph()
        p2.text = "   " + desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_after = Pt(8)

    # ==================== SLIDE 7: Dịch Vụ Excel, Media & Security ====================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "TÍNH NĂNG 4: DỊCH VỤ EXCEL, MEDIA & BẢO MẬT", "功能 4: Excel 服务、媒体与权限安全")

    s7_cards = [
        {
            "title_vi": "📑 Xuất / Nhập Excel",
            "title_zh": "Excel 导入/导出",
            "sub": "Dịch vụ excel.ts",
            "items": [
                ("Xuất báo cáo đơn mua ra (.xlsx)", "导出采购报告 (.xlsx)"),
                ("Nhập hàng loạt mã vật tư", "批量导入物料编码数据"),
                ("Tự động chuẩn hóa dữ liệu", "自动标准化数据格式")
            ]
        },
        {
            "title_vi": "🖼️ Quản Lý Hình Ảnh",
            "title_zh": "图片管理服务",
            "sub": "Dịch vụ imageService.ts",
            "items": [
                ("Upload ảnh vật tư lên Storage", "上传物料图片至 Storage"),
                ("Xem trực tiếp ảnh trong Modal", "详情弹窗直接预览图片"),
                ("Tối ưu dung lượng & tốc độ", "优化图片体积与加载")
            ]
        },
        {
            "title_vi": "🔒 Phân Quyền & Bảo Mật",
            "title_zh": "权限与安全控制",
            "sub": "Supabase Auth & RLS",
            "items": [
                ("Phân quyền Role User & Admin", "严格划分 User 与 Admin 权限"),
                ("Trang AdminUserManagement", "管理员用户管理页面"),
                ("Bảo mật Row Level Security", "PostgreSQL 行级安全策略")
            ]
        }
    ]

    for i, sc in enumerate(s7_cards):
        left = Inches(0.8 + i * 4.04)
        create_card(slide7, left, Inches(1.8), Inches(3.64), Inches(4.9))
        
        tb = slide7.shapes.add_textbox(left + Inches(0.15), Inches(1.95), Inches(3.34), Inches(4.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = sc["title_vi"] + "\n" + sc["title_zh"]
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.space_after = Pt(2)

        p_s = tf.add_paragraph()
        p_s.text = sc["sub"]
        p_s.font.size = Pt(10.5)
        p_s.font.color.rgb = TEXT_MUTED
        p_s.space_after = Pt(10)

        for vi, zh in sc["items"]:
            pi_vi = tf.add_paragraph()
            pi_vi.text = "• " + vi
            pi_vi.font.size = Pt(11)
            pi_vi.font.bold = True
            pi_vi.font.color.rgb = TEXT_DARK
            
            pi_zh = tf.add_paragraph()
            pi_zh.text = "  " + zh
            pi_zh.font.size = Pt(10)
            pi_zh.font.color.rgb = TEXT_MUTED
            pi_zh.space_after = Pt(8)

    # ==================== SLIDE 8: Chứng Chỉ Bảo Mật Database ====================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "CHỨNG CHỈ BẢO MẬT & TUÂN THỦ CƠ SỞ DỮ LIỆU (SUPABASE / POSTGRESQL)", "数据库安全与合规认证 (SOC 2, ISO 27001, HIPAA, GDPR)")

    cert_cards = [
        {
            "badge": "🛡️ SOC 2 Type II Certified",
            "title_vi": "Chứng Nhận SOC 2 Type II",
            "title_zh": "SOC 2 Type II 审计认证",
            "desc_vi": "Đạt kiểm toán an ninh độc lập hàng năm về Bảo mật, Tính sẵn sàng và Bảo mật thông tin.",
            "desc_zh": "通过年度独立安全审计，涵盖安全性、可用性与机密性。"
        },
        {
            "badge": "🔒 ISO 27001 Certified",
            "title_vi": "Tiêu Chuẩn ISO 27001",
            "title_zh": "ISO 27001 国际认证",
            "desc_vi": "Chứng nhận quốc tế về Hệ thống quản lý an toàn thông tin (ISMS) cấp độ doanh nghiệp.",
            "desc_zh": "企业级信息安全管理体系 (ISMS) 国际标准认证。"
        },
        {
            "badge": "⚕️ HIPAA Compliant",
            "title_vi": "Tuân Thủ HIPAA",
            "title_zh": "HIPAA 合规认证",
            "desc_vi": "Đáp ứng đầy đủ quy chuẩn bảo mật dữ liệu y tế & doanh nghiệp khắt khe (BAA Ready).",
            "desc_zh": "满足严格的高标准医疗与企业数据安全规范。"
        },
        {
            "badge": "🇪🇺 GDPR Compliant",
            "title_vi": "Tuân Thủ GDPR",
            "title_zh": "GDPR 隐私合规",
            "desc_vi": "Bảo đảm quyền riêng tư dữ liệu cá nhân theo luật Châu Âu, hỗ trợ Data Residency tại Châu Á.",
            "desc_zh": "遵循欧洲数据隐私保护法规，支持亚洲数据驻留。"
        },
        {
            "badge": "💳 PCI-DSS Level 1",
            "title_vi": "Bảo Mật Thanh Toán",
            "title_zh": "PCI-DSS 1级安全",
            "desc_vi": "Hạ tầng lưu trữ đạt tiêu chuẩn an toàn dữ liệu thanh toán thẻ PCI-DSS Cấp độ 1.",
            "desc_zh": "存储基础设施达到支付卡数据安全最高一级标准。"
        },
        {
            "badge": "🔑 Engine Security",
            "title_vi": "Mã Hóa & RLS Engine",
            "title_zh": "引擎级加密与 RLS",
            "desc_vi": "Mã hóa AES-256 (At Rest), TLS 1.3 (In Transit), Row Level Security (RLS) trên Postgres.",
            "desc_zh": "AES-256 (静态加密), TLS 1.3 (传输加密) 及 Postgres RLS 行级安全。"
        }
    ]

    cert_grid_w = Inches(3.64)
    cert_grid_h = Inches(2.3)

    cert_coords = [
        (Inches(0.8), Inches(1.8)),
        (Inches(4.84), Inches(1.8)),
        (Inches(8.88), Inches(1.8)),
        (Inches(0.8), Inches(4.35)),
        (Inches(4.84), Inches(4.35)),
        (Inches(8.88), Inches(4.35))
    ]

    for i, cc in enumerate(cert_cards):
        left, top = cert_coords[i]
        create_card(slide8, left, top, cert_grid_w, cert_grid_h, bg_color=CARD_BG, border_color=PRIMARY)
        
        tb = slide8.shapes.add_textbox(left + Inches(0.12), top + Inches(0.1), cert_grid_w - Inches(0.24), cert_grid_h - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_badge = tf.paragraphs[0]
        p_badge.text = cc["badge"].upper()
        p_badge.font.size = Pt(10)
        p_badge.font.bold = True
        p_badge.font.color.rgb = PRIMARY
        p_badge.space_after = Pt(2)

        p_t = tf.add_paragraph()
        p_t.text = cc["title_vi"] + " • " + cc["title_zh"]
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK
        p_t.space_after = Pt(4)

        p_dvi = tf.add_paragraph()
        p_dvi.text = cc["desc_vi"]
        p_dvi.font.size = Pt(10)
        p_dvi.font.color.rgb = TEXT_DARK

        p_dzh = tf.add_paragraph()
        p_dzh.text = cc["desc_zh"]
        p_dzh.font.size = Pt(9.5)
        p_dzh.font.color.rgb = TEXT_MUTED

    # ==================== SLIDE 9: Tổng Kết ====================
    slide9 = prs.slides.add_slide(blank_layout)
    bg9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg9.fill.solid()
    bg9.fill.fore_color.rgb = DARK_BG
    bg9.line.fill.background()

    tb9 = slide9.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.333), Inches(5.2))
    tf9 = tb9.text_frame
    tf9.word_wrap = True

    p = tf9.paragraphs[0]
    p.text = "TỔNG KẾT & ĐỊNH HƯỚNG MỞ RỘNG • 总结与扩展方向"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(20)

    conclusions = [
        ("✨ Quy trình khép kín / 闭环流程", "Tối ưu từ đề xuất, duyệt đơn, tạo PO đến nghiệm thu nhập kho.\n从申请、审批、生成PO到验收入库全流程优化。"),
        ("⚡ Hiệu năng vượt trội / 卓越性能", "React 19 + Vite 6 cho trải nghiệm người dùng tức thì.\nReact 19 + Vite 6 带来即时响应体验。"),
        ("🛡️ Bảo mật tiêu chuẩn / 国际标准安全", "Hạ tầng Supabase Postgres đạt chứng nhận SOC 2, ISO 27001, HIPAA & GDPR.\nSupabase Postgres 基础设施通过 SOC 2, ISO 27001, HIPAA 和 GDPR 认证。"),
        ("🌐 Song ngữ Việt - Trung / 越中双语", "Phù hợp hoàn hảo cho doanh nghiệp & nhà máy sản xuất liên doanh.\n完美契合合资企业与生产工厂需求。"),
        ("🚀 Khả năng mở rộng / 扩展能力", "Tích hợp AI gợi ý mua hàng & kết nối hệ thống ERP (SAP).\n集成 AI 采购建议并对接 SAP ERP 系统。")
    ]

    for title, desc in conclusions:
        pi_t = tf9.add_paragraph()
        pi_t.text = title
        pi_t.font.size = Pt(14)
        pi_t.font.bold = True
        pi_t.font.color.rgb = PRIMARY_LIGHT
        
        pi_d = tf9.add_paragraph()
        pi_d.text = "   " + desc
        pi_d.font.size = Pt(11)
        pi_d.font.color.rgb = TEXT_LIGHT
        pi_d.space_after = Pt(10)

    # Output file
    output_path = os.path.join(os.getcwd(), "Bao_Cao_Tinh_Nang_Du_An_Mua_Hang_Bilingual.pptx")
    try:
        prs.save(output_path)
    except PermissionError:
        output_path = os.path.join(os.getcwd(), "Bao_Cao_Tinh_Nang_Du_An_Mua_Hang_Bilingual_v2.pptx")
        prs.save(output_path)

    if sys.stdout.encoding.lower() != 'utf-8':
        try:
            sys.stdout.reconfigure(encoding='utf-8')
        except Exception:
            pass
    print(f"SUCCESS: Created Bilingual PowerPoint file at: {output_path}")

if __name__ == "__main__":
    build_presentation()
